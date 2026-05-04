import hashlib
from pathlib import Path
from xml.etree import ElementTree
from xml.sax.saxutils import escape as xml_escape

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt


CANVAS_PX = (1152, 648)
SLIDE_W = Inches(16)
SLIDE_H = Inches(9)
EMU_PER_INCH = 914400
PX_PER_INCH = CANVAS_PX[0] / 16
SVG_NS = 'http://www.w3.org/2000/svg'
STRICT_SVG_ALLOWED_TAGS = {
    f'{{{SVG_NS}}}svg',
    f'{{{SVG_NS}}}g',
    f'{{{SVG_NS}}}rect',
    f'{{{SVG_NS}}}text',
}

PALETTES = [
    {'bg': '#F6F2EA', 'ink': '#171C24', 'muted': '#5B6570', 'accent': '#B94624', 'panel': '#EFE6D6', 'line': '#D8C8B2'},
    {'bg': '#EDF4F2', 'ink': '#10211D', 'muted': '#53645F', 'accent': '#147C72', 'panel': '#DCEAE4', 'line': '#B9D1C8'},
    {'bg': '#F5F2F8', 'ink': '#211B2B', 'muted': '#625B6E', 'accent': '#6E4AA8', 'panel': '#E6E0F0', 'line': '#CFC3E3'},
    {'bg': '#F4F6F8', 'ink': '#152033', 'muted': '#5D6675', 'accent': '#2E6FBB', 'panel': '#E1E8F0', 'line': '#C4D0DD'},
]

RECIPE_TO_LAYOUT = {
    'ppt.hero_signal': 'cover_signal',
    'ppt.compare_zones': 'multi_zone_compare',
    'ppt.timeline_rail': 'timeline_band',
    'ppt.judgement_ladder': 'judgement_ladder',
    'ppt.ring_cross': 'ring_cross',
    'ppt.summary_peak': 'summary_peak',
}

SUPPORTED_LAYOUTS = {
    'cover_signal',
    'multi_zone_compare',
    'timeline_band',
    'judgement_ladder',
    'ring_cross',
    'summary_peak',
}


def safe_text(value, fallback: str = '') -> str:
    text = str(value or '').strip()
    return text or fallback


def safe_list(value):
    return value if isinstance(value, list) else []


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip('#')
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def emu_to_px(value) -> float:
    return float(value) / EMU_PER_INCH * PX_PER_INCH


def shape_rect(left, top, width, height) -> dict:
    x = emu_to_px(left)
    y = emu_to_px(top)
    w = emu_to_px(width)
    h = emu_to_px(height)
    return {
        'left': round(x, 2),
        'top': round(y, 2),
        'width': round(w, 2),
        'height': round(h, 2),
        'right': round(x + w, 2),
        'bottom': round(y + h, 2),
    }


def file_sha256(file: Path) -> str:
    if not file.exists():
        return ''
    digest = hashlib.sha256()
    with file.open('rb') as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b''):
            digest.update(chunk)
    return digest.hexdigest()


def visible_text(value, fallback: str = '') -> str:
    if value is None:
        return fallback
    if isinstance(value, str):
        return safe_text(value, fallback)
    if isinstance(value, (int, float)):
        return safe_text(value, fallback)
    if isinstance(value, list):
        return safe_text('\n'.join(visible_text(item) for item in value if visible_text(item)), fallback)
    if isinstance(value, dict):
        for key in ('editable_text', 'text', 'body', 'title', 'core_sentence'):
            text = visible_text(value.get(key))
            if text:
                return text
        return fallback
    return fallback


def visible_text_lines(value, limit: int = 5) -> list[str]:
    lines = []
    source = value if isinstance(value, list) else [value]
    for item in source:
        text = visible_text(item)
        for line in text.splitlines():
            line = safe_text(line)
            if line:
                lines.append(line)
    return lines[:limit]


def layout_family(slide_data: dict) -> str:
    nested = slide_data.get('visual_presentation') if isinstance(slide_data.get('visual_presentation'), dict) else {}
    layout = safe_text(slide_data.get('layout_family') or nested.get('layout_family'))
    if not layout:
        layout = RECIPE_TO_LAYOUT.get(safe_text(slide_data.get('render_recipe_id')), '')
    return layout if layout in SUPPORTED_LAYOUTS else 'multi_zone_compare'


def layout_writer_id(layout: str) -> str:
    return f'{layout}_native_writer'


def shape_text(shape: dict) -> str:
    return visible_text(shape.get('editable_text') or shape.get('text'))


def slide_title(slide_data, index: int) -> str:
    for shape in safe_list(slide_data.get('_editable_native_shapes')):
        if isinstance(shape, dict) and safe_text(shape.get('role')) == 'title':
            text = shape_text(shape)
            if text:
                return text
    return visible_text(slide_data.get('title'), f'Slide {index}')


def slide_core_sentence(slide_data) -> str:
    for shape in safe_list(slide_data.get('_editable_native_shapes')):
        if isinstance(shape, dict) and safe_text(shape.get('role')) in {'core_sentence', 'subtitle'}:
            text = shape_text(shape)
            if text:
                return text
    return visible_text(slide_data.get('core_sentence'))


def slide_points(slide_data) -> list[str]:
    plan_points = []
    for shape in safe_list(slide_data.get('_editable_native_shapes')):
        if not isinstance(shape, dict):
            continue
        if safe_text(shape.get('role')) not in {'point_text', 'bullet', 'body', 'content', 'content_stack'}:
            continue
        plan_points.extend(visible_text_lines(shape_text(shape), limit=6))
    if plan_points:
        return plan_points[:5]
    points = visible_text_lines(slide_data.get('page_core_content'), limit=5)
    if points:
        return points
    points = visible_text_lines(slide_data.get('evidence_points'), limit=5)
    if points:
        return points
    return [visible_text(slide_data.get('core_sentence'), 'Key message')]


def assert_authoring_text_capacity(text: str, limit: int, shape_id: str) -> str:
    text = safe_text(text)
    if len(text) > limit:
        raise ValueError(
            'native PPT authoring text exceeds deterministic writer capacity: '
            f'{shape_id} has {len(text)} chars, limit {limit}'
        )
    return text


def svg_rect(x, y, width, height, fill, intent, shape_id='', stroke='none', rx=0):
    attrs = [
        f'x="{x}"',
        f'y="{y}"',
        f'width="{width}"',
        f'height="{height}"',
        f'fill="{xml_escape(fill)}"',
        f'stroke="{xml_escape(stroke)}"',
        f'data-intent="{xml_escape(intent)}"',
    ]
    if shape_id:
        attrs.append(f'id="{xml_escape(shape_id)}"')
    if rx:
        attrs.append(f'rx="{rx}"')
    return f'    <rect {" ".join(attrs)} />'


def svg_text(x, y, text, size, fill, intent, shape_id='', weight='400'):
    checked_text = assert_authoring_text_capacity(text, 160, shape_id or intent)
    attrs = [
        f'x="{x}"',
        f'y="{y}"',
        f'font-size="{size}"',
        f'font-family="Aptos, Arial, sans-serif"',
        f'font-weight="{weight}"',
        f'fill="{xml_escape(fill)}"',
        f'data-intent="{xml_escape(intent)}"',
    ]
    if shape_id:
        attrs.append(f'id="{xml_escape(shape_id)}"')
    return f'    <text {" ".join(attrs)}>{xml_escape(checked_text)}</text>'


def strict_svg_preflight(svg_file: Path) -> dict:
    try:
        root = ElementTree.parse(svg_file).getroot()
    except ElementTree.ParseError as exc:
        raise ValueError(f'SVG IR preflight failed for {svg_file}: {exc}') from exc
    if root.tag != f'{{{SVG_NS}}}svg':
        raise ValueError(f'SVG IR preflight failed for {svg_file}: root element must be svg')
    if root.attrib.get('data-redcube-ir') != 'redcube_svg_ir_v1':
        raise ValueError(f'SVG IR preflight failed for {svg_file}: missing redcube_svg_ir_v1 marker')
    counts = {'text': 0, 'rect': 0, 'group': 0}
    for element in root.iter():
        if element.tag not in STRICT_SVG_ALLOWED_TAGS:
            raise ValueError(f'SVG IR preflight failed for {svg_file}: unsupported tag {element.tag}')
        local = element.tag.split('}', 1)[-1]
        if local == 'text':
            counts['text'] += 1
        elif local == 'rect':
            counts['rect'] += 1
        elif local == 'g':
            counts['group'] += 1
    if counts['text'] < 1 or counts['rect'] < 1 or counts['group'] < 1:
        raise ValueError(f'SVG IR preflight failed for {svg_file}: text/rect/group intent is required')
    return {
        'status': 'pass',
        'strict': True,
        'allowed_tags': ['svg', 'g', 'rect', 'text'],
        'intent_counts': counts,
    }


def build_slide_svg_ir(slide_data, index: int, total: int, palette: dict, repaired: bool) -> str:
    slide_id = safe_text(slide_data.get('slide_id'), f'S{index:02d}')
    layout = layout_family(slide_data)
    points = slide_points(slide_data)
    lines = [
        f'<svg xmlns="{SVG_NS}" width="{CANVAS_PX[0]}" height="{CANVAS_PX[1]}" viewBox="0 0 {CANVAS_PX[0]} {CANVAS_PX[1]}" data-redcube-ir="redcube_svg_ir_v1" data-slide-id="{xml_escape(slide_id)}" data-layout-writer="{xml_escape(layout_writer_id(layout))}">',
        '  <g id="slide-root" data-intent="group:slide">',
        svg_rect(0, 0, CANVAS_PX[0], CANVAS_PX[1], palette['bg'], 'rect:background', f'{slide_id}-background'),
        svg_text(72, 76, slide_title(slide_data, index), 34 if layout != 'cover_signal' else 48, palette['ink'], 'text:title', f'{slide_id}-title', '700'),
    ]
    if repaired:
        lines.append(svg_rect(1032, 36, 52, 12, palette['accent'], 'rect:repair_marker', f'{slide_id}-repair-marker'))
    if layout == 'cover_signal':
        lines.extend([
            svg_rect(72, 148, 720, 6, palette['accent'], 'rect:hero_rule', f'{slide_id}-hero-rule'),
            svg_text(72, 214, slide_core_sentence(slide_data), 25, palette['muted'], 'text:core_sentence', f'{slide_id}-core-sentence', '600'),
        ])
        x = 92
        for point_index, point in enumerate(points[:3], 1):
            lines.extend([
                f'    <g id="{xml_escape(f"{slide_id}-signal-{point_index}")}" data-intent="group:signal_point">',
                svg_rect(x, 408, 288, 92, palette['panel'], 'rect:signal_panel', f'{slide_id}-signal-{point_index}-panel', palette['line'], 16),
                svg_text(x + 24, 462, point, 18, palette['ink'], 'text:point_text', f'{slide_id}-signal-{point_index}-text', '600'),
                '    </g>',
            ])
            x += 330
    elif layout == 'timeline_band':
        lines.append(svg_text(72, 142, slide_core_sentence(slide_data), 18, palette['muted'], 'text:core_sentence', f'{slide_id}-core-sentence'))
        lines.append(svg_rect(104, 326, 944, 8, palette['accent'], 'rect:timeline_rail', f'{slide_id}-timeline-rail'))
        for point_index, point in enumerate(points[:4], 1):
            x = 92 + ((point_index - 1) * 252)
            lines.extend([
                f'    <g id="{xml_escape(f"{slide_id}-milestone-{point_index}")}" data-intent="group:timeline_point">',
                svg_rect(x, 254, 198, 118, palette['panel'], 'rect:timeline_panel', f'{slide_id}-milestone-{point_index}-panel', palette['line'], 14),
                svg_text(x + 18, 296, f'{point_index:02d}', 22, palette['accent'], 'text:point_index', f'{slide_id}-milestone-{point_index}-index', '700'),
                svg_text(x + 18, 340, point, 16, palette['ink'], 'text:point_text', f'{slide_id}-milestone-{point_index}-text', '600'),
                '    </g>',
            ])
    elif layout == 'judgement_ladder':
        lines.append(svg_text(72, 142, slide_core_sentence(slide_data), 18, palette['muted'], 'text:core_sentence', f'{slide_id}-core-sentence'))
        for point_index, point in enumerate(points[:4], 1):
            x = 118 + ((point_index - 1) * 220)
            y = 428 - ((point_index - 1) * 58)
            lines.extend([
                f'    <g id="{xml_escape(f"{slide_id}-gate-{point_index}")}" data-intent="group:judgement_gate">',
                svg_rect(x, y, 210, 112, palette['panel'], 'rect:judgement_step', f'{slide_id}-gate-{point_index}-panel', palette['line'], 12),
                svg_text(x + 18, y + 38, f'Gate {point_index}', 18, palette['accent'], 'text:point_index', f'{slide_id}-gate-{point_index}-index', '700'),
                svg_text(x + 18, y + 78, point, 15, palette['ink'], 'text:point_text', f'{slide_id}-gate-{point_index}-text', '600'),
                '    </g>',
            ])
    elif layout == 'ring_cross':
        lines.append(svg_text(72, 142, slide_core_sentence(slide_data), 18, palette['muted'], 'text:core_sentence', f'{slide_id}-core-sentence'))
        lines.append(svg_rect(462, 262, 228, 118, palette['accent'], 'rect:center_hub', f'{slide_id}-center-hub', palette['accent'], 24))
        lines.append(svg_text(496, 326, 'Core', 24, '#FFFFFF', 'text:center_hub', f'{slide_id}-center-text', '700'))
        positions = [(122, 214), (758, 214), (758, 426), (122, 426)]
        for point_index, point in enumerate(points[:4], 1):
            x, y = positions[point_index - 1]
            lines.extend([
                f'    <g id="{xml_escape(f"{slide_id}-axis-{point_index}")}" data-intent="group:ring_cross_axis">',
                svg_rect(x, y, 286, 116, palette['panel'], 'rect:axis_panel', f'{slide_id}-axis-{point_index}-panel', palette['line'], 14),
                svg_text(x + 22, y + 48, point, 16, palette['ink'], 'text:point_text', f'{slide_id}-axis-{point_index}-text', '600'),
                '    </g>',
            ])
    elif layout == 'summary_peak':
        lines.append(svg_text(72, 146, slide_core_sentence(slide_data), 20, palette['muted'], 'text:core_sentence', f'{slide_id}-core-sentence', '600'))
        lines.append(svg_rect(96, 214, 960, 120, palette['accent'], 'rect:summary_peak', f'{slide_id}-summary-peak', palette['accent'], 26))
        lines.append(svg_text(132, 284, points[0] if points else slide_core_sentence(slide_data), 24, '#FFFFFF', 'text:summary_peak', f'{slide_id}-peak-text', '700'))
        for point_index, point in enumerate(points[1:4], 1):
            x = 96 + ((point_index - 1) * 326)
            lines.extend([
                svg_rect(x, 394, 296, 96, palette['panel'], 'rect:takeaway_panel', f'{slide_id}-takeaway-{point_index}-panel', palette['line'], 16),
                svg_text(x + 22, 452, point, 16, palette['ink'], 'text:point_text', f'{slide_id}-takeaway-{point_index}-text', '600'),
            ])
    else:
        lines.append(svg_text(72, 142, slide_core_sentence(slide_data), 18, palette['muted'], 'text:core_sentence', f'{slide_id}-core-sentence'))
        columns = [(84, 220, 456, 242), (612, 220, 456, 242)]
        for point_index, point in enumerate(points[:4], 1):
            column = 0 if point_index <= 2 else 1
            row = 0 if point_index in {1, 3} else 1
            x, y, width, height = columns[column]
            panel_y = y + row * 124
            lines.extend([
                f'    <g id="{xml_escape(f"{slide_id}-zone-{point_index}")}" data-intent="group:compare_zone">',
                svg_rect(x, panel_y, width, 102, palette['panel'], 'rect:compare_panel', f'{slide_id}-zone-{point_index}-panel', palette['line'], 14),
                svg_text(x + 24, panel_y + 58, point, 16, palette['ink'], 'text:point_text', f'{slide_id}-zone-{point_index}-text', '600'),
                '    </g>',
            ])
    footer = f'{slide_id} / {index} of {total}'
    lines.extend([
        svg_text(72, 614, footer, 15, palette['muted'], 'text:page_number', f'{slide_id}-footer'),
        '  </g>',
        '</svg>',
    ])
    return '\n'.join(lines) + '\n'


def render_slide_svg_ir(slide_data, index: int, total: int, svg_dir: Path, palette: dict, repaired: bool) -> dict:
    slide_id = safe_text(slide_data.get('slide_id'), f'S{index:02d}')
    svg_file = svg_dir / f'{index:02d}-{slide_id}.svg'
    svg_file.parent.mkdir(parents=True, exist_ok=True)
    svg_file.write_text(build_slide_svg_ir(slide_data, index, total, palette, repaired), encoding='utf-8')
    preflight = strict_svg_preflight(svg_file)
    return {
        'file': str(svg_file),
        'sha256': file_sha256(svg_file),
        'preflight': preflight,
    }


class SlideBuilder:
    def __init__(self, slide, palette: dict, slide_id: str):
        self.slide = slide
        self.palette = palette
        self.slide_id = slide_id
        self.native_shapes = []

    def record_shape(self, shape_id, kind, left, top, width, height, role, text='', font_size=0, quality_role='content'):
        self.native_shapes.append({
            'shape_id': shape_id,
            'kind': kind,
            'role': role,
            'quality_role': quality_role,
            'bounds': shape_rect(left, top, width, height),
            'text': safe_text(text),
            'font_size': font_size,
        })

    def rect(self, shape_id, left, top, width, height, fill_key, role, quality_role='decorative', line_key=None, radius=True):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = self.slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(self.palette[fill_key] if fill_key in self.palette else fill_key)
        line_color = self.palette.get(line_key or fill_key, self.palette.get(fill_key, fill_key))
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(1)
        shape.name = shape_id
        self.record_shape(shape_id, 'rounded_rect' if radius else 'rect', left, top, width, height, role, quality_role=quality_role)
        return shape

    def oval(self, shape_id, left, top, width, height, fill_key, role, quality_role='decorative', line_key=None):
        shape = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(self.palette[fill_key] if fill_key in self.palette else fill_key)
        line_color = self.palette.get(line_key or fill_key, self.palette.get(fill_key, fill_key))
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(1)
        shape.name = shape_id
        self.record_shape(shape_id, 'oval', left, top, width, height, role, quality_role=quality_role)
        return shape

    def line(self, shape_id, begin_x, begin_y, end_x, end_y, color_key='line', width=1.2, role='connector', quality_role='decorative'):
        begin_x = int(begin_x)
        begin_y = int(begin_y)
        end_x = int(end_x)
        end_y = int(end_y)
        shape = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, begin_x, begin_y, end_x, end_y)
        shape.name = shape_id
        shape.line.color.rgb = rgb(self.palette[color_key] if color_key in self.palette else color_key)
        shape.line.width = Pt(width)
        left = min(begin_x, end_x)
        top = min(begin_y, end_y)
        box_width = abs(end_x - begin_x)
        box_height = abs(end_y - begin_y)
        self.record_shape(shape_id, 'line', left, top, box_width, box_height, role, quality_role=quality_role)
        return shape

    def text(self, shape_id, left, top, width, height, text, size, color_key, role, bold=False, quality_role='content'):
        text = assert_authoring_text_capacity(text, 220, shape_id)
        box = self.slide.shapes.add_textbox(left, top, width, height)
        box.name = shape_id
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.05)
        frame.margin_right = Inches(0.05)
        frame.margin_top = Inches(0.03)
        frame.margin_bottom = Inches(0.03)
        paragraph = frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(self.palette[color_key] if color_key in self.palette else color_key)
        self.record_shape(shape_id, 'text_box', left, top, width, height, role, text, size, quality_role)
        return box


def add_background(ctx: SlideBuilder):
    ctx.rect(f'{ctx.slide_id}-background', 0, 0, SLIDE_W, SLIDE_H, 'bg', 'background', radius=False)
    ctx.rect(f'{ctx.slide_id}-top-rule', Inches(1.0), Inches(0.36), Inches(1.85), Inches(0.08), 'accent', 'accent_rule', radius=False)
    ctx.line(f'{ctx.slide_id}-left-rail', Inches(0.62), Inches(0.8), Inches(0.62), Inches(7.7), 'line', 1.1, 'layout_rail')
    ctx.oval(f'{ctx.slide_id}-corner-dot', Inches(14.42), Inches(0.48), Inches(0.18), Inches(0.18), 'accent', 'accent_dot')
    ctx.rect(f'{ctx.slide_id}-corner-chip', Inches(14.7), Inches(0.54), Inches(0.34), Inches(0.08), 'line', 'accent_chip', radius=False)


def add_footer(ctx: SlideBuilder, index: int, total: int):
    ctx.text(
        f'{ctx.slide_id}-footer',
        Inches(1.0),
        Inches(8.27),
        Inches(13.7),
        Inches(0.3),
        f'{ctx.slide_id} / {index} of {total}',
        10,
        'muted',
        'page_number',
        quality_role='decorative',
    )


def add_title_band(ctx: SlideBuilder, title: str, core: str, title_size=30):
    ctx.text(f'{ctx.slide_id}-title', Inches(1.0), Inches(0.58), Inches(13.9), Inches(0.72), title, title_size, 'ink', 'title', True)
    if core:
        ctx.text(f'{ctx.slide_id}-core-sentence', Inches(1.0), Inches(1.38), Inches(13.5), Inches(0.58), core, 15, 'muted', 'core_sentence')


def write_cover_signal(ctx: SlideBuilder, slide_data, index: int, total: int, repaired: bool):
    add_background(ctx)
    title = slide_title(slide_data, index)
    core = slide_core_sentence(slide_data)
    points = slide_points(slide_data)
    ctx.rect(f'{ctx.slide_id}-hero-block', Inches(0.95), Inches(1.38), Inches(10.25), Inches(2.2), 'panel', 'hero_panel', 'content', 'line')
    ctx.text(f'{ctx.slide_id}-title', Inches(1.28), Inches(1.68), Inches(9.45), Inches(0.92), title, 34, 'ink', 'title', True)
    ctx.text(f'{ctx.slide_id}-core-sentence', Inches(1.28), Inches(2.68), Inches(9.4), Inches(0.54), core, 17, 'muted', 'core_sentence')
    ctx.rect(f'{ctx.slide_id}-accent-slab', Inches(12.0), Inches(0.95), Inches(2.55), Inches(6.6), 'accent', 'hero_accent', 'decorative', radius=False)
    ctx.line(f'{ctx.slide_id}-hero-vector', Inches(10.8), Inches(2.18), Inches(12.0), Inches(2.18), 'accent', 1.6, 'hero_vector')
    ctx.oval(f'{ctx.slide_id}-hero-node', Inches(11.24), Inches(2.04), Inches(0.28), Inches(0.28), 'bg', 'hero_node', 'decorative', 'accent')
    for point_index, point in enumerate(points[:3], 1):
        x = Inches(1.0 + ((point_index - 1) * 4.55))
        ctx.rect(f'{ctx.slide_id}-signal-{point_index}-panel', x, Inches(5.35), Inches(3.95), Inches(1.0), 'panel', 'signal_panel', 'content', 'line')
        ctx.oval(f'{ctx.slide_id}-signal-{point_index}-dot', x + Inches(0.18), Inches(5.12), Inches(0.18), Inches(0.18), 'accent', 'signal_dot')
        ctx.text(f'{ctx.slide_id}-signal-{point_index}-text', x + Inches(0.24), Inches(5.58), Inches(3.4), Inches(0.48), point, 14, 'ink', 'point_text', True)
    add_footer(ctx, index, total)


def write_multi_zone_compare(ctx: SlideBuilder, slide_data, index: int, total: int, repaired: bool):
    add_background(ctx)
    add_title_band(ctx, slide_title(slide_data, index), slide_core_sentence(slide_data))
    points = slide_points(slide_data)
    zones = [
        (Inches(0.95), Inches(2.25), Inches(6.65), Inches(2.15)),
        (Inches(8.05), Inches(2.25), Inches(6.65), Inches(2.15)),
        (Inches(0.95), Inches(4.72), Inches(6.65), Inches(1.78)),
        (Inches(8.05), Inches(4.72), Inches(6.65), Inches(1.78)),
    ]
    for point_index, point in enumerate(points[:4], 1):
        left, top, width, height = zones[point_index - 1]
        ctx.rect(f'{ctx.slide_id}-zone-{point_index}-panel', left, top, width, height, 'panel', 'compare_panel', 'content', 'line')
        ctx.rect(f'{ctx.slide_id}-zone-{point_index}-tab', left + Inches(0.28), top + Inches(0.18), Inches(1.05), Inches(0.12), 'accent', 'zone_tab', radius=False)
        ctx.text(f'{ctx.slide_id}-zone-{point_index}-index', left + Inches(0.28), top + Inches(0.24), Inches(0.6), Inches(0.38), f'{point_index:02d}', 14, 'accent', 'point_index', True)
        ctx.text(f'{ctx.slide_id}-zone-{point_index}-text', left + Inches(0.28), top + Inches(0.74), width - Inches(0.58), height - Inches(0.9), point, 14, 'ink', 'point_text')
    ctx.rect(f'{ctx.slide_id}-bridge-rule', Inches(7.78), Inches(2.38), Inches(0.08), Inches(4.0), 'accent', 'compare_divider', radius=False)
    ctx.line(f'{ctx.slide_id}-horizontal-bridge', Inches(1.2), Inches(4.55), Inches(14.55), Inches(4.55), 'line', 1.1, 'compare_bridge')
    add_footer(ctx, index, total)


def write_timeline_band(ctx: SlideBuilder, slide_data, index: int, total: int, repaired: bool):
    add_background(ctx)
    add_title_band(ctx, slide_title(slide_data, index), slide_core_sentence(slide_data))
    points = slide_points(slide_data)
    ctx.rect(f'{ctx.slide_id}-timeline-rail', Inches(1.3), Inches(4.02), Inches(13.0), Inches(0.08), 'accent', 'timeline_rail', radius=False)
    for point_index, point in enumerate(points[:4], 1):
        left = Inches(1.0 + ((point_index - 1) * 3.65))
        top = Inches(2.55 if point_index % 2 else 4.45)
        center_x = left + Inches(1.52)
        ctx.line(f'{ctx.slide_id}-milestone-{point_index}-stem', center_x, Inches(4.06), center_x, top + Inches(0.08), 'line', 1.1, 'timeline_stem')
        ctx.oval(f'{ctx.slide_id}-milestone-{point_index}-node', center_x - Inches(0.11), Inches(3.93), Inches(0.22), Inches(0.22), 'accent', 'timeline_node')
        ctx.rect(f'{ctx.slide_id}-milestone-{point_index}-panel', left, top, Inches(3.05), Inches(1.15), 'panel', 'timeline_panel', 'content', 'line')
        ctx.text(f'{ctx.slide_id}-milestone-{point_index}-index', left + Inches(0.22), top + Inches(0.18), Inches(0.65), Inches(0.32), f'{point_index:02d}', 13, 'accent', 'point_index', True)
        ctx.text(f'{ctx.slide_id}-milestone-{point_index}-text', left + Inches(0.22), top + Inches(0.56), Inches(2.55), Inches(0.46), point, 12.5, 'ink', 'point_text')
    add_footer(ctx, index, total)


def write_judgement_ladder(ctx: SlideBuilder, slide_data, index: int, total: int, repaired: bool):
    add_background(ctx)
    add_title_band(ctx, slide_title(slide_data, index), slide_core_sentence(slide_data))
    points = slide_points(slide_data)
    ctx.line(f'{ctx.slide_id}-ladder-spine', Inches(0.95), Inches(6.62), Inches(12.85), Inches(3.84), 'line', 1.5, 'ladder_spine')
    for point_index, point in enumerate(points[:4], 1):
        left = Inches(1.08 + ((point_index - 1) * 3.18))
        top = Inches(5.8 - ((point_index - 1) * 0.72))
        height = Inches(0.92 + ((point_index - 1) * 0.32))
        ctx.oval(f'{ctx.slide_id}-gate-{point_index}-node', left - Inches(0.18), top + Inches(0.16), Inches(0.22), Inches(0.22), 'accent', 'gate_node')
        ctx.rect(f'{ctx.slide_id}-gate-{point_index}-panel', left, top, Inches(2.76), height, 'panel', 'judgement_step', 'content', 'line')
        ctx.text(f'{ctx.slide_id}-gate-{point_index}-index', left + Inches(0.2), top + Inches(0.16), Inches(1.3), Inches(0.3), f'Gate {point_index}', 13, 'accent', 'point_index', True)
        ctx.text(f'{ctx.slide_id}-gate-{point_index}-text', left + Inches(0.2), top + Inches(0.54), Inches(2.3), height - Inches(0.68), point, 12.5, 'ink', 'point_text')
    add_footer(ctx, index, total)


def write_ring_cross(ctx: SlideBuilder, slide_data, index: int, total: int, repaired: bool):
    add_background(ctx)
    add_title_band(ctx, slide_title(slide_data, index), slide_core_sentence(slide_data))
    points = slide_points(slide_data)
    ctx.rect(f'{ctx.slide_id}-center-hub', Inches(6.25), Inches(3.36), Inches(3.18), Inches(1.05), 'accent', 'center_hub', 'content')
    ctx.text(f'{ctx.slide_id}-center-text', Inches(6.72), Inches(3.64), Inches(2.25), Inches(0.4), 'Core', 18, '#FFFFFF', 'center_hub', True)
    positions = [
        (Inches(1.02), Inches(2.42), Inches(4.35), Inches(1.3)),
        (Inches(10.55), Inches(2.42), Inches(4.35), Inches(1.3)),
        (Inches(10.55), Inches(5.06), Inches(4.35), Inches(1.3)),
        (Inches(1.02), Inches(5.06), Inches(4.35), Inches(1.3)),
    ]
    hub_points = [
        (Inches(6.25), Inches(3.62)),
        (Inches(9.43), Inches(3.62)),
        (Inches(9.43), Inches(4.12)),
        (Inches(6.25), Inches(4.12)),
    ]
    for point_index, point in enumerate(points[:4], 1):
        left, top, width, height = positions[point_index - 1]
        end_x, end_y = hub_points[point_index - 1]
        start_x = left + (width if point_index in {1, 4} else 0)
        start_y = top + (height / 2)
        ctx.line(f'{ctx.slide_id}-axis-{point_index}-connector', start_x, start_y, end_x, end_y, 'line', 1.2, 'axis_connector')
        ctx.rect(f'{ctx.slide_id}-axis-{point_index}-panel', left, top, width, height, 'panel', 'axis_panel', 'content', 'line')
        ctx.oval(f'{ctx.slide_id}-axis-{point_index}-node', left + Inches(0.16), top + Inches(0.16), Inches(0.18), Inches(0.18), 'accent', 'axis_node')
        ctx.text(f'{ctx.slide_id}-axis-{point_index}-text', left + Inches(0.24), top + Inches(0.32), width - Inches(0.5), Inches(0.62), point, 13.5, 'ink', 'point_text', True)
    add_footer(ctx, index, total)


def write_summary_peak(ctx: SlideBuilder, slide_data, index: int, total: int, repaired: bool):
    add_background(ctx)
    add_title_band(ctx, slide_title(slide_data, index), slide_core_sentence(slide_data), 31)
    points = slide_points(slide_data)
    primary = points[0] if points else slide_core_sentence(slide_data)
    ctx.rect(f'{ctx.slide_id}-summary-peak', Inches(1.0), Inches(2.38), Inches(13.7), Inches(1.35), 'accent', 'summary_peak', 'content')
    ctx.rect(f'{ctx.slide_id}-peak-shadow', Inches(1.26), Inches(3.82), Inches(13.18), Inches(0.12), 'line', 'peak_shadow', radius=False)
    ctx.text(f'{ctx.slide_id}-peak-text', Inches(1.48), Inches(2.75), Inches(12.5), Inches(0.54), primary, 18, '#FFFFFF', 'summary_peak', True)
    for point_index, point in enumerate(points[1:4] or points[:3], 1):
        left = Inches(1.0 + ((point_index - 1) * 4.63))
        ctx.oval(f'{ctx.slide_id}-takeaway-{point_index}-node', left + Inches(0.2), Inches(4.28), Inches(0.22), Inches(0.22), 'accent', 'takeaway_node')
        ctx.rect(f'{ctx.slide_id}-takeaway-{point_index}-panel', left, Inches(4.55), Inches(4.05), Inches(1.25), 'panel', 'takeaway_panel', 'content', 'line')
        ctx.text(f'{ctx.slide_id}-takeaway-{point_index}-text', left + Inches(0.26), Inches(4.92), Inches(3.5), Inches(0.48), point, 13.5, 'ink', 'point_text', True)
    add_footer(ctx, index, total)


LAYOUT_WRITERS = {
    'cover_signal': write_cover_signal,
    'multi_zone_compare': write_multi_zone_compare,
    'timeline_band': write_timeline_band,
    'judgement_ladder': write_judgement_ladder,
    'ring_cross': write_ring_cross,
    'summary_peak': write_summary_peak,
}


def build_deck(slides, output_pptx: Path, svg_ir_dir: Path, repaired_slide_ids, evaluate_native_slide_quality):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]
    manifest_slides = []
    for index, slide_data in enumerate(slides, 1):
        palette = PALETTES[(index - 1) % len(PALETTES)]
        slide = prs.slides.add_slide(blank_layout)
        slide_id = safe_text(slide_data.get('slide_id'), f'S{index:02d}')
        repaired = slide_id in repaired_slide_ids
        layout = layout_family(slide_data)
        writer = LAYOUT_WRITERS[layout]
        try:
            svg_ir = render_slide_svg_ir(slide_data, index, len(slides), svg_ir_dir, palette, repaired)
        except ValueError as exc:
            raise RuntimeError(str(exc)) from exc
        ctx = SlideBuilder(slide, palette, slide_id)
        writer(ctx, slide_data, index, len(slides), repaired)
        if repaired:
            ctx.rect(f'{slide_id}-repair-marker', Inches(14.18), Inches(0.38), Inches(0.72), Inches(0.12), 'accent', 'repair_marker', radius=False)
        points = slide_points(slide_data)
        quality = evaluate_native_slide_quality(ctx.native_shapes, len(points))
        manifest_slides.append({
            'slide_id': slide_id,
            'title': slide_title(slide_data, index),
            'layout_family': layout,
            'layout_writer': layout_writer_id(layout),
            'title_font_size': 34 if layout == 'cover_signal' else 30,
            'shape_count': len(slide.shapes),
            'text_box_count': sum(1 for shape in slide.shapes if getattr(shape, 'has_text_frame', False)),
            'redcube_svg_ir_file': svg_ir['file'],
            'redcube_svg_ir_sha256': svg_ir['sha256'],
            'redcube_svg_ir_preflight': svg_ir['preflight'],
            'redcube_svg_ir': {
                'file': svg_ir['file'],
                'sha256': svg_ir['sha256'],
                'preflight': svg_ir['preflight'],
            },
            'preview_screenshot_file': '',
            'preview_screenshot_sha256': '',
            'preview_screenshot_dimensions': None,
            'render_proof_source': '',
            'synthetic_preview': False,
            'native_shapes': ctx.native_shapes,
            'checks': quality['checks'],
            'metrics': quality['metrics'],
            'issues': quality['issues'],
            'repaired': repaired,
        })
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))
    return manifest_slides
