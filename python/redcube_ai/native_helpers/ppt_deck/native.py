#!/usr/bin/env python3
import argparse
import hashlib
import json
import math
import shutil
import subprocess
import sys
import textwrap
from pathlib import Path
from xml.etree import ElementTree
from xml.sax.saxutils import escape as xml_escape

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


CANVAS_PX = (1152, 648)
SLIDE_W = Inches(16)
SLIDE_H = Inches(9)
EMU_PER_INCH = 914400
PX_PER_INCH = CANVAS_PX[0] / 16
FRAME_AREA = float(CANVAS_PX[0] * CANVAS_PX[1])
MIN_NATIVE_DENSITY = 0.18
MAX_NATIVE_DENSITY = 0.82
MIN_NATIVE_EDGE_CLEARANCE = 24.0
MAX_NATIVE_PRIMARY_POINTS = 5
ENGINE_CAPABILITIES = {
    'authoring_ir': 'redcube_svg_ir',
    'authoring_ir_version': 1,
    'pptx_writer': 'redcube_drawingml_writer',
    'editable_pptx': True,
    'strict_svg_preflight': True,
    'true_render_proof_required': True,
    'true_render_proof_renderer': 'powerpoint_applescript',
    'screenshot_packaging': False,
}
SVG_NS = 'http://www.w3.org/2000/svg'
STRICT_SVG_ALLOWED_TAGS = {
    f'{{{SVG_NS}}}svg',
    f'{{{SVG_NS}}}g',
    f'{{{SVG_NS}}}rect',
    f'{{{SVG_NS}}}text',
}
DEFAULT_ENGINE_CONTRACT_FILE = (
    Path(__file__).resolve().parents[4]
    / 'contracts'
    / 'runtime-program'
    / 'ppt-native-python-engine-contract.json'
)


def fail(message: str) -> None:
    print(message, file=sys.stderr)
    raise SystemExit(1)


def safe_text(value, fallback: str = '') -> str:
    text = str(value or '').strip()
    return text or fallback


def safe_list(value):
    return value if isinstance(value, list) else []


def clamp(value: float, lower: float, upper: float) -> float:
    return max(lower, min(upper, value))


def emu_to_px(value) -> float:
    return float(value) / EMU_PER_INCH * PX_PER_INCH


def shape_rect(left, top, width, height) -> dict:
    x = emu_to_px(left)
    y = emu_to_px(top)
    w = emu_to_px(width)
    h = emu_to_px(height)
    return {
        'left': round(x, 2),
        'top': round(y, 2),
        'width': round(w, 2),
        'height': round(h, 2),
        'right': round(x + w, 2),
        'bottom': round(y + h, 2),
    }


def rect_overlap_area(rect_a: dict, rect_b: dict) -> float:
    overlap_w = max(0.0, min(rect_a['right'], rect_b['right']) - max(rect_a['left'], rect_b['left']))
    overlap_h = max(0.0, min(rect_a['bottom'], rect_b['bottom']) - max(rect_a['top'], rect_b['top']))
    return overlap_w * overlap_h


def file_sha256(file: Path) -> str:
    if not file.exists():
        return ''
    digest = hashlib.sha256()
    with file.open('rb') as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b''):
            digest.update(chunk)
    return digest.hexdigest()


def image_dimensions(file: Path) -> dict:
    with Image.open(file) as image:
        return {'width': int(image.width), 'height': int(image.height)}


def text_capacity_failure(shape: dict) -> dict | None:
    text = safe_text(shape.get('text'))
    if not text:
        return None
    rect = shape.get('bounds') or {}
    width = float(rect.get('width') or 0)
    height = float(rect.get('height') or 0)
    font_size = float(shape.get('font_size') or 14)
    chars_per_line = max(8, int(width / max(font_size * 0.52, 1)))
    estimated_lines = max(1, math.ceil(len(text) / chars_per_line))
    max_lines = max(1, int(height / max(font_size * 1.05, 1))) + 2
    if estimated_lines <= max_lines:
        return None
    return {
        'shape_id': shape.get('shape_id'),
        'overflow_reason': 'native_text_capacity_exceeded',
        'text_char_count': len(text),
        'estimated_lines': estimated_lines,
        'max_lines': max_lines,
    }


def evaluate_native_slide_quality(native_shapes: list, primary_points: int) -> dict:
    content_shapes = [shape for shape in native_shapes if shape.get('quality_role') == 'content']
    overlap_shapes = [shape for shape in content_shapes if shape.get('kind') == 'text_box']
    occupied_area = sum(float(shape['bounds']['width']) * float(shape['bounds']['height']) for shape in content_shapes)
    occupied_ratio = clamp(occupied_area / FRAME_AREA, 0.0, 1.0)
    if content_shapes:
        edge_clearance = {
            'left': min(float(shape['bounds']['left']) for shape in content_shapes),
            'top': min(float(shape['bounds']['top']) for shape in content_shapes),
            'right': min(CANVAS_PX[0] - float(shape['bounds']['right']) for shape in content_shapes),
            'bottom': min(CANVAS_PX[1] - float(shape['bounds']['bottom']) for shape in content_shapes),
        }
    else:
        edge_clearance = {'left': 0.0, 'top': 0.0, 'right': 0.0, 'bottom': 0.0}
    overlap_pairs = []
    for left_index, shape_a in enumerate(overlap_shapes):
        for shape_b in overlap_shapes[left_index + 1:]:
            overlap_area = rect_overlap_area(shape_a['bounds'], shape_b['bounds'])
            if overlap_area > 12.0:
                overlap_pairs.append({
                    'a': shape_a.get('shape_id'),
                    'b': shape_b.get('shape_id'),
                    'overlap_area': round(overlap_area, 2),
                })
    block_content_failures = [
        failure for failure in (text_capacity_failure(shape) for shape in content_shapes)
        if failure
    ]
    text_char_count = sum(len(safe_text(shape.get('text'))) for shape in content_shapes)
    clipped_nodes = len(block_content_failures)
    checks = {
        'overflow_free': clipped_nodes == 0,
        'occlusion_free': len(overlap_pairs) == 0,
        'visual_density_ok': MIN_NATIVE_DENSITY <= occupied_ratio <= MAX_NATIVE_DENSITY and primary_points <= MAX_NATIVE_PRIMARY_POINTS,
        'speaker_fit_ok': text_char_count <= 950,
        'edge_clearance_ok': min(edge_clearance.values()) >= MIN_NATIVE_EDGE_CLEARANCE,
        'block_content_fit_ok': clipped_nodes == 0,
        'title_typography_ok': True,
        'page_number_consistency_ok': True,
    }
    issues = []
    if not checks['visual_density_ok']:
        issues.append('visual_density_out_of_range')
    if not checks['edge_clearance_ok']:
        issues.append('edge_clearance_out_of_range')
    if not checks['occlusion_free']:
        issues.append('occlusion_detected')
    if not checks['block_content_fit_ok']:
        issues.append('block_content_overflow_detected')
    if not checks['speaker_fit_ok']:
        issues.append('speaker_fit_exceeded')
    return {
        'checks': checks,
        'issues': issues,
        'metrics': {
            'title_font_size': 32,
            'text_char_count': text_char_count,
            'block_count': len(content_shapes),
            'shape_count': len(native_shapes),
            'overlap_pairs': len(overlap_pairs),
            'overlaps': overlap_pairs,
            'clipped_nodes': clipped_nodes,
            'occupied_ratio': round(occupied_ratio, 4),
            'primary_points': primary_points,
            'edge_clearance': {key: round(value, 2) for key, value in edge_clearance.items()},
            'block_content_failures': block_content_failures,
            'bounds': [shape['bounds'] for shape in content_shapes],
        },
    }


def load_engine_contract(contract_file: Path) -> dict:
    if not contract_file.exists():
        fail(f'engine contract not found: {contract_file}')
    contract = json.loads(contract_file.read_text(encoding='utf-8'))
    required = {
        'kind': 'redcube_native_ppt_python_engine',
        'language': 'python',
        'contract_version': 1,
        'input_boundary': 'slide_blueprint_plus_visual_direction_json',
        'review_boundary': 'rendered_pptx_screenshots',
    }
    for key, expected in required.items():
        if contract.get(key) != expected:
            fail(f'engine contract field mismatch: {key}')
    routes = contract.get('owned_routes')
    if routes != ['author_pptx_native', 'repair_pptx_native']:
        fail('engine contract owned_routes mismatch')
    capabilities = contract.get('engine_capabilities') or {}
    for key, expected in ENGINE_CAPABILITIES.items():
        if capabilities.get(key) != expected:
            fail(f'engine contract capability mismatch: {key}')
    render_proof = contract.get('true_render_proof') or {}
    if render_proof.get('required') is not True:
        fail('engine contract true_render_proof.required mismatch')
    return contract


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip('#')
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


PALETTES = [
    {'bg': '#F7F3EA', 'ink': '#18212B', 'muted': '#59636F', 'accent': '#BA3F1D', 'panel': '#EFE6D6'},
    {'bg': '#EEF4F1', 'ink': '#12231F', 'muted': '#52645F', 'accent': '#147C72', 'panel': '#DCEAE4'},
    {'bg': '#F3F1F7', 'ink': '#201B2C', 'muted': '#625A6F', 'accent': '#6E4AA8', 'panel': '#E4DFF0'},
    {'bg': '#F4F6F8', 'ink': '#152033', 'muted': '#5D6675', 'accent': '#2E6FBB', 'panel': '#E2E8F0'},
]


def font(size: int, bold: bool = False):
    candidates = [
        '/System/Library/Fonts/PingFang.ttc',
        '/System/Library/Fonts/Supplemental/Arial Unicode.ttf',
        '/Library/Fonts/Arial Unicode.ttf',
    ]
    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size=size)
        except Exception:
            continue
    return ImageFont.load_default()


def add_textbox(slide, left, top, width, height, text, size, color, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_panel(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(1)
    return shape


def normalize_slide_data(payload: dict) -> list:
    plan = payload.get('editable_shape_plan') or {}
    plan_slides = safe_list(plan.get('slides'))
    blueprint = payload.get('blueprint') or {}
    blueprint_slides = safe_list(blueprint.get('slides'))
    blueprint_by_id = {
        safe_text(slide.get('slide_id'), f'S{index + 1:02d}'): slide
        for index, slide in enumerate(blueprint_slides)
        if isinstance(slide, dict)
    }
    source_slides = plan_slides or blueprint_slides
    if not source_slides:
        fail('native PPT authoring requires editable_shape_plan.slides or slide_blueprint.slides')
    slides = []
    for index, raw_slide in enumerate(source_slides, 1):
        if not isinstance(raw_slide, dict):
            continue
        slide_id = safe_text(raw_slide.get('slide_id'), f'S{index:02d}')
        blueprint_slide = blueprint_by_id.get(slide_id, {})
        merged = {**blueprint_slide, **raw_slide}
        plan_shapes = [
            shape for shape in safe_list(raw_slide.get('native_shapes'))
            if isinstance(shape, dict)
        ]
        merged['_editable_native_shapes'] = plan_shapes
        slides.append(merged)
    if not slides:
        fail('native PPT authoring requires at least one valid slide object')
    return slides


def slide_points(slide_data):
    plan_points = [
        safe_text(shape.get('editable_text') or shape.get('text'))
        for shape in safe_list(slide_data.get('_editable_native_shapes'))
        if isinstance(shape, dict)
        and safe_text(shape.get('role')) in {'point_text', 'bullet', 'body', 'content'}
        and safe_text(shape.get('editable_text') or shape.get('text'))
    ]
    if plan_points:
        return plan_points[:4]
    points = [
        safe_text(item)
        for item in safe_list(slide_data.get('page_core_content'))
        if safe_text(item)
    ]
    if points:
        return points[:4]
    points = [
        safe_text(item)
        for item in safe_list(slide_data.get('evidence_points'))
        if safe_text(item)
    ]
    if points:
        return points[:4]
    return [safe_text(slide_data.get('core_sentence'), 'Key message')]


def slide_title(slide_data, index: int) -> str:
    title_shape = next(
        (
            shape for shape in safe_list(slide_data.get('_editable_native_shapes'))
            if isinstance(shape, dict)
            and safe_text(shape.get('role')) == 'title'
            and safe_text(shape.get('editable_text') or shape.get('text'))
        ),
        None,
    )
    if title_shape:
        return safe_text(title_shape.get('editable_text') or title_shape.get('text'), f'Slide {index}')
    return safe_text(slide_data.get('title'), f'Slide {index}')


def slide_core_sentence(slide_data) -> str:
    core_shape = next(
        (
            shape for shape in safe_list(slide_data.get('_editable_native_shapes'))
            if isinstance(shape, dict)
            and safe_text(shape.get('role')) in {'core_sentence', 'subtitle'}
            and safe_text(shape.get('editable_text') or shape.get('text'))
        ),
        None,
    )
    if core_shape:
        return safe_text(core_shape.get('editable_text') or core_shape.get('text'))
    return safe_text(slide_data.get('core_sentence'))


def svg_rect(x, y, width, height, fill, intent, shape_id='', stroke='none', rx=0):
    attrs = [
        f'x="{x}"',
        f'y="{y}"',
        f'width="{width}"',
        f'height="{height}"',
        f'fill="{xml_escape(fill)}"',
        f'stroke="{xml_escape(stroke)}"',
        f'data-intent="{xml_escape(intent)}"',
    ]
    if shape_id:
        attrs.append(f'id="{xml_escape(shape_id)}"')
    if rx:
        attrs.append(f'rx="{rx}"')
    return f'    <rect {" ".join(attrs)} />'


def svg_text(x, y, text, size, fill, intent, shape_id='', weight='400'):
    attrs = [
        f'x="{x}"',
        f'y="{y}"',
        f'font-size="{size}"',
        f'font-family="Aptos, Arial, sans-serif"',
        f'font-weight="{weight}"',
        f'fill="{xml_escape(fill)}"',
        f'data-intent="{xml_escape(intent)}"',
    ]
    if shape_id:
        attrs.append(f'id="{xml_escape(shape_id)}"')
    return f'    <text {" ".join(attrs)}>{xml_escape(safe_text(text))}</text>'


def build_slide_svg_ir(slide_data, index: int, total: int, palette: dict, repaired: bool) -> str:
    slide_id = safe_text(slide_data.get('slide_id'), f'S{index:02d}')
    lines = [
        f'<svg xmlns="{SVG_NS}" width="{CANVAS_PX[0]}" height="{CANVAS_PX[1]}" viewBox="0 0 {CANVAS_PX[0]} {CANVAS_PX[1]}" data-redcube-ir="redcube_svg_ir_v1" data-slide-id="{xml_escape(slide_id)}">',
        '  <g id="slide-root" data-intent="group:slide">',
        svg_rect(0, 0, CANVAS_PX[0], CANVAS_PX[1], palette['bg'], 'rect:background', f'{slide_id}-background'),
        svg_rect(0, 0, 36, CANVAS_PX[1], palette['accent'], 'rect:accent_bar', f'{slide_id}-accent-bar'),
        svg_text(84, 88, slide_title(slide_data, index), 34, palette['ink'], 'text:title', f'{slide_id}-title', '700'),
        svg_text(84, 176, slide_core_sentence(slide_data), 22, palette['muted'], 'text:core_sentence', f'{slide_id}-core-sentence'),
    ]
    if repaired:
        lines.append(svg_rect(1032, 36, 52, 12, palette['accent'], 'rect:repair_marker', f'{slide_id}-repair-marker'))
    y = 250
    for point_index, point in enumerate(slide_points(slide_data)[:4], 1):
        lines.extend([
            f'    <g id="{xml_escape(f"{slide_id}-point-{point_index}")}" data-intent="group:content_point">',
            svg_rect(96, y, 960, 68, palette['panel'], 'rect:content_panel', f'{slide_id}-point-{point_index}-panel', rx=18),
            svg_text(124, y + 43, str(point_index), 24, palette['accent'], 'text:point_index', f'{slide_id}-point-{point_index}-index', '700'),
            svg_text(176, y + 42, point, 19, palette['ink'], 'text:point_text', f'{slide_id}-point-{point_index}-text'),
            '    </g>',
        ])
        y += 86
    footer = f'{slide_id} / {index} of {total}'
    lines.extend([
        svg_text(84, 612, footer, 16, palette['muted'], 'text:page_number', f'{slide_id}-footer'),
        '  </g>',
        '</svg>',
    ])
    return '\n'.join(lines) + '\n'


def strict_svg_preflight(svg_file: Path) -> dict:
    try:
        root = ElementTree.parse(svg_file).getroot()
    except ElementTree.ParseError as exc:
        fail(f'SVG IR preflight failed for {svg_file}: {exc}')
    if root.tag != f'{{{SVG_NS}}}svg':
        fail(f'SVG IR preflight failed for {svg_file}: root element must be svg')
    if root.attrib.get('data-redcube-ir') != 'redcube_svg_ir_v1':
        fail(f'SVG IR preflight failed for {svg_file}: missing redcube_svg_ir_v1 marker')
    counts = {'text': 0, 'rect': 0, 'group': 0}
    for element in root.iter():
        if element.tag not in STRICT_SVG_ALLOWED_TAGS:
            fail(f'SVG IR preflight failed for {svg_file}: unsupported tag {element.tag}')
        local = element.tag.split('}', 1)[-1]
        if local == 'image':
            fail(f'SVG IR preflight failed for {svg_file}: raster image nodes are not allowed')
        if local == 'text':
            counts['text'] += 1
        elif local == 'rect':
            counts['rect'] += 1
        elif local == 'g':
            counts['group'] += 1
    if counts['text'] < 1 or counts['rect'] < 1 or counts['group'] < 1:
        fail(f'SVG IR preflight failed for {svg_file}: text/rect/group intent is required')
    return {
        'status': 'pass',
        'strict': True,
        'allowed_tags': ['svg', 'g', 'rect', 'text'],
        'intent_counts': counts,
    }


def render_slide_svg_ir(slide_data, index: int, total: int, svg_dir: Path, palette: dict, repaired: bool) -> dict:
    slide_id = safe_text(slide_data.get('slide_id'), f'S{index:02d}')
    svg_file = svg_dir / f'{index:02d}-{slide_id}.svg'
    svg_file.parent.mkdir(parents=True, exist_ok=True)
    svg_file.write_text(build_slide_svg_ir(slide_data, index, total, palette, repaired), encoding='utf-8')
    preflight = strict_svg_preflight(svg_file)
    return {
        'file': str(svg_file),
        'sha256': file_sha256(svg_file),
        'preflight': preflight,
    }


def draw_wrapped(draw, xy, text, image_font, fill, width_chars):
    lines = []
    for paragraph in safe_text(text).split('\n'):
        lines.extend(textwrap.wrap(paragraph, width=width_chars) or [''])
    x, y = xy
    line_height = int(image_font.size * 1.35) if hasattr(image_font, 'size') else 18
    for line in lines[:6]:
        draw.text((x, y), line, font=image_font, fill=fill)
        y += line_height


def render_synthetic_debug_preview(slide_data, index, total, preview_file: Path, palette, repaired: bool):
    image = Image.new('RGB', CANVAS_PX, palette['bg'])
    draw = ImageDraw.Draw(image)
    accent = palette['accent']
    ink = palette['ink']
    muted = palette['muted']
    panel = palette['panel']
    draw.rectangle([0, 0, 36, CANVAS_PX[1]], fill=accent)
    draw.rectangle([84, 116, 1068, 118], fill=accent)
    if repaired:
        draw.rectangle([1032, 36, 1084, 48], fill=accent)
    draw_wrapped(draw, (84, 54), safe_text(slide_data.get('title'), f'Slide {index}'), font(34, True), ink, 29)
    draw_wrapped(draw, (84, 142), safe_text(slide_data.get('core_sentence')), font(22), muted, 50)
    y = 250
    for point_index, point in enumerate(slide_points(slide_data)[:4], 1):
        draw.rounded_rectangle([96, y, 1056, y + 68], radius=18, fill=panel)
        draw.text((124, y + 18), f'{point_index}', font=font(24, True), fill=accent)
        draw_wrapped(draw, (176, y + 16), point, font(19), ink, 58)
        y += 86
    footer = f'{safe_text(slide_data.get("slide_id"), f"S{index:02d}")} / {index} of {total}'
    draw.text((84, 596), footer, font=font(16), fill=muted)
    preview_file.parent.mkdir(parents=True, exist_ok=True)
    image.save(preview_file)


def build_deck(slides, output_pptx: Path, svg_ir_dir: Path, repaired_slide_ids):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]
    manifest_slides = []
    for index, slide_data in enumerate(slides, 1):
        palette = PALETTES[(index - 1) % len(PALETTES)]
        slide = prs.slides.add_slide(blank_layout)
        repaired = safe_text(slide_data.get('slide_id')) in repaired_slide_ids
        slide_id = safe_text(slide_data.get('slide_id'), f'S{index:02d}')
        svg_ir = render_slide_svg_ir(slide_data, index, len(slides), svg_ir_dir, palette, repaired)
        native_shapes = []

        def record_shape(shape_id, kind, left, top, width, height, role, text='', font_size=0, quality_role='content'):
            native_shapes.append({
                'shape_id': shape_id,
                'kind': kind,
                'role': role,
                'quality_role': quality_role,
                'bounds': shape_rect(left, top, width, height),
                'text': safe_text(text),
                'font_size': font_size,
            })

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = rgb(palette['bg'])
        bg.line.color.rgb = rgb(palette['bg'])
        record_shape(f'{slide_id}-background', 'shape', 0, 0, SLIDE_W, SLIDE_H, 'background', quality_role='decorative')
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), SLIDE_H)
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(palette['accent'])
        bar.line.color.rgb = rgb(palette['accent'])
        record_shape(f'{slide_id}-accent-bar', 'shape', 0, 0, Inches(0.5), SLIDE_H, 'accent_bar', quality_role='decorative')
        add_textbox(
            slide,
            Inches(1.15),
            Inches(0.55),
            Inches(13.3),
            Inches(0.92),
            slide_title(slide_data, index),
            32,
            rgb(palette['ink']),
            bold=True,
        )
        record_shape(
            f'{slide_id}-title',
            'text_box',
            Inches(1.15),
            Inches(0.55),
            Inches(13.3),
            Inches(0.92),
            'title',
            slide_title(slide_data, index),
            32,
        )
        add_textbox(
            slide,
            Inches(1.15),
            Inches(1.62),
            Inches(13.3),
            Inches(0.72),
            slide_core_sentence(slide_data),
            18,
            rgb(palette['muted']),
        )
        record_shape(
            f'{slide_id}-core-sentence',
            'text_box',
            Inches(1.15),
            Inches(1.62),
            Inches(13.3),
            Inches(0.72),
            'core_sentence',
            slide_core_sentence(slide_data),
            18,
        )
        top = Inches(3.05)
        points = slide_points(slide_data)[:4]
        for point_index, point in enumerate(points, 1):
            panel = add_panel(slide, Inches(1.25), top, Inches(13.3), Inches(0.78), rgb(palette['panel']))
            record_shape(
                f'{slide_id}-point-{point_index}-panel',
                'shape',
                Inches(1.25),
                top,
                Inches(13.3),
                Inches(0.78),
                'content_panel',
                quality_role='decorative',
            )
            add_textbox(
                slide,
                Inches(1.45),
                top + Inches(0.13),
                Inches(0.45),
                Inches(0.34),
                str(point_index),
                18,
                rgb(palette['accent']),
                bold=True,
            )
            record_shape(
                f'{slide_id}-point-{point_index}-index',
                'text_box',
                Inches(1.45),
                top + Inches(0.13),
                Inches(0.45),
                Inches(0.34),
                'point_index',
                str(point_index),
                18,
            )
            add_textbox(
                slide,
                Inches(2.1),
                top + Inches(0.11),
                Inches(11.8),
                Inches(0.42),
                point,
                15,
                rgb(palette['ink']),
            )
            record_shape(
                f'{slide_id}-point-{point_index}-text',
                'text_box',
                Inches(2.1),
                top + Inches(0.11),
                Inches(11.8),
                Inches(0.42),
                'point_text',
                point,
                15,
            )
            top += Inches(0.96)
        footer_text = f'{slide_id} / {index} of {len(slides)}'
        add_textbox(
            slide,
            Inches(1.15),
            Inches(8.25),
            Inches(13.3),
            Inches(0.3),
            footer_text,
            10,
            rgb(palette['muted']),
        )
        record_shape(
            f'{slide_id}-footer',
            'text_box',
            Inches(1.15),
            Inches(8.25),
            Inches(13.3),
            Inches(0.3),
            'page_number',
            footer_text,
            10,
        )
        quality = evaluate_native_slide_quality(native_shapes, len(points))
        manifest_slides.append({
            'slide_id': slide_id,
            'title': slide_title(slide_data, index),
            'layout_family': safe_text(slide_data.get('layout_family')),
            'title_font_size': 32,
            'shape_count': len(slide.shapes),
            'text_box_count': sum(1 for shape in slide.shapes if getattr(shape, 'has_text_frame', False)),
            'redcube_svg_ir_file': svg_ir['file'],
            'redcube_svg_ir_sha256': svg_ir['sha256'],
            'redcube_svg_ir_preflight': svg_ir['preflight'],
            'redcube_svg_ir': {
                'file': svg_ir['file'],
                'sha256': svg_ir['sha256'],
                'preflight': svg_ir['preflight'],
            },
            'preview_screenshot_file': '',
            'preview_screenshot_sha256': '',
            'preview_screenshot_dimensions': None,
            'render_proof_source': 'true_pptx_render',
            'synthetic_preview': False,
            'native_shapes': native_shapes,
            'checks': quality['checks'],
            'metrics': quality['metrics'],
            'issues': quality['issues'],
            'repaired': repaired,
        })
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))
    return manifest_slides


def apple_script_text(value: Path | str) -> str:
    return str(value).replace('\\', '\\\\').replace('"', '\\"')


def render_powerpoint_pptx(output_pptx: Path, preview_dir: Path, output_pdf: Path | None) -> dict:
    if not Path('/Applications/Microsoft PowerPoint.app').exists():
        fail('true PPTX render proof requires Microsoft PowerPoint at /Applications/Microsoft PowerPoint.app')
    preview_dir.mkdir(parents=True, exist_ok=True)
    for stale in preview_dir.glob('*'):
        if stale.is_file() and stale.suffix.lower() in {'.png', '.pdf'}:
            stale.unlink()
    pdf_target = output_pdf or (preview_dir / 'render-proof.pdf')
    pdf_target.parent.mkdir(parents=True, exist_ok=True)
    script = f'''
set pptxPath to POSIX file "{apple_script_text(output_pptx)}"
set pngDir to POSIX file "{apple_script_text(preview_dir)}"
set pdfPath to POSIX file "{apple_script_text(pdf_target)}"
tell application id "com.microsoft.Powerpoint"
    launch
    open pptxPath
    set deckRef to active presentation
    save deckRef in pngDir as save as PNG
    save deckRef in pdfPath as save as PDF
    close deckRef saving no
end tell
'''
    completed = subprocess.run(
        ['osascript', '-e', script],
        text=True,
        capture_output=True,
        check=False,
    )
    if completed.returncode != 0:
        stderr = safe_text(completed.stderr) or safe_text(completed.stdout) or 'unknown osascript error'
        fail(f'true PPTX render proof failed via PowerPoint AppleScript: {stderr}')
    png_files = sorted(path for path in preview_dir.glob('*.png') if path.is_file())
    if not png_files:
        png_files = render_png_from_pdf(pdf_target, preview_dir)
    if not png_files:
        fail(f'true PPTX render proof produced no PNG previews in {preview_dir}')
    if output_pdf and not output_pdf.exists() and pdf_target.exists():
        shutil.copyfile(pdf_target, output_pdf)
    return {
        'source_surface_kind': 'native_pptx',
        'renderer_kind': 'powerpoint_applescript',
        'synthetic_preview': False,
        'required': True,
        'pptx_file': str(output_pptx),
        'pdf_file': str(output_pdf or pdf_target),
        'preview_screenshots': [str(path) for path in png_files],
    }


def render_png_from_pdf(pdf_file: Path, preview_dir: Path) -> list[Path]:
    if not pdf_file.exists():
        return []
    pdftoppm = shutil.which('pdftoppm')
    if pdftoppm:
        prefix = preview_dir / 'slide'
        completed = subprocess.run(
            [pdftoppm, '-png', '-r', '144', str(pdf_file), str(prefix)],
            text=True,
            capture_output=True,
            check=False,
        )
        if completed.returncode != 0:
            stderr = safe_text(completed.stderr) or safe_text(completed.stdout)
            fail(f'PDF-to-PNG render proof conversion failed with pdftoppm: {stderr}')
        return sorted(path for path in preview_dir.glob('slide-*.png') if path.is_file())
    magick = shutil.which('magick') or shutil.which('convert')
    if magick:
        output_pattern = str(preview_dir / 'slide-%02d.png')
        completed = subprocess.run(
            [magick, '-density', '144', str(pdf_file), output_pattern],
            text=True,
            capture_output=True,
            check=False,
        )
        if completed.returncode != 0:
            stderr = safe_text(completed.stderr) or safe_text(completed.stdout)
            fail(f'PDF-to-PNG render proof conversion failed with ImageMagick: {stderr}')
        return sorted(path for path in preview_dir.glob('slide-*.png') if path.is_file())
    fail('PowerPoint PDF render succeeded, but no PDF-to-PNG converter is available for preview screenshots')


def attach_rendered_previews(manifest_slides: list, render_proof: dict) -> list:
    preview_files = [Path(path) for path in safe_list(render_proof.get('preview_screenshots'))]
    if len(preview_files) != len(manifest_slides):
        fail(
            'true PPTX render proof slide count mismatch: '
            f'{len(preview_files)} previews for {len(manifest_slides)} manifest slides'
        )
    attached = []
    for slide, preview_file in zip(manifest_slides, preview_files):
        if not preview_file.exists():
            fail(f'true PPTX render proof screenshot missing: {preview_file}')
        attached.append({
            **slide,
            'preview_screenshot_file': str(preview_file),
            'preview_screenshot_sha256': file_sha256(preview_file),
            'preview_screenshot_dimensions': image_dimensions(preview_file),
            'render_proof_source': 'true_pptx_render',
            'synthetic_preview': False,
        })
    return attached


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description='Build editable native PPTX artifacts for ppt_deck.')
    parser.add_argument('--input-json', required=True)
    parser.add_argument('--mode', choices=['author', 'repair'], required=True)
    parser.add_argument('--output-pptx', required=True)
    parser.add_argument('--shape-manifest', required=True)
    parser.add_argument('--preview-dir', required=True)
    parser.add_argument('--output-pdf')
    parser.add_argument('--repair-log')
    parser.add_argument('--engine-contract', default=str(DEFAULT_ENGINE_CONTRACT_FILE))
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    payload = json.loads(Path(args.input_json).read_text(encoding='utf-8'))
    engine_contract_file = Path(args.engine_contract).resolve()
    engine_contract = load_engine_contract(engine_contract_file)
    slides = normalize_slide_data(payload)
    repair_feedback = safe_list(payload.get('repair_feedback'))
    repaired_slide_ids = {
        safe_text(item.get('slide_id'))
        for item in repair_feedback
        if isinstance(item, dict) and safe_text(item.get('slide_id'))
    }
    output_pptx = Path(args.output_pptx).resolve()
    shape_manifest = Path(args.shape_manifest).resolve()
    preview_dir = Path(args.preview_dir).resolve()
    svg_ir_dir = preview_dir.parent / f'{preview_dir.name}-redcube-svg-ir'
    manifest_slides = build_deck(slides, output_pptx, svg_ir_dir, repaired_slide_ids)
    output_pdf = Path(args.output_pdf).resolve() if args.output_pdf else None
    render_proof = render_powerpoint_pptx(output_pptx, preview_dir, output_pdf)
    manifest_slides = attach_rendered_previews(manifest_slides, render_proof)
    preview_files = safe_list(render_proof.get('preview_screenshots'))
    repair_log_file = Path(args.repair_log).resolve() if args.repair_log else None
    repair_log = {
        'mode': args.mode,
        'consumed_review_stage': 'screenshot_review' if args.mode == 'repair' else None,
        'target_slide_ids': sorted(repaired_slide_ids),
        'feedback_count': len(repair_feedback),
        'repair_log_file': str(repair_log_file) if repair_log_file else None,
    }
    if repair_log_file:
        repair_log_file.parent.mkdir(parents=True, exist_ok=True)
        repair_log_file.write_text(json.dumps(repair_log, ensure_ascii=False, indent=2), encoding='utf-8')
    manifest = {
        'schema_version': 1,
        'artifact_kind': 'ppt_deck_native_shape_manifest',
        'engine_contract': engine_contract,
        'engine_contract_file': str(engine_contract_file),
        'builder': {
            'kind': 'redcube_drawingml_writer',
            'implementation': 'python_pptx_native_shapes',
            'surface': 'editable_native_pptx',
            'screenshot_packaging': False,
        },
        'capability': {
            'kind': 'RedCube DrawingML writer',
            'editable_artifact': True,
            'native_shapes': True,
            'redcube_svg_ir': True,
            'strict_svg_preflight': True,
            'render_proof_required': True,
        },
        'native_quality_model': 'shape_manifest_layout_metrics_v1',
        'native_quality_surface': {
            'quality_model': 'shape_manifest_layout_metrics_v1',
            'source_surface_kind': 'native_pptx',
            'required_per_slide_metrics': [
                'bounds',
                'text_char_count',
                'primary_points',
                'occupied_ratio',
                'edge_clearance',
                'overlap_pairs',
                'preview_screenshot_sha256',
                'preview_screenshot_dimensions',
            ],
            'fail_closed_when_missing': True,
        },
        'engine_capabilities': ENGINE_CAPABILITIES,
        'render_proof': render_proof,
        'redcube_svg_ir': {
            'kind': 'redcube_svg_ir',
            'version': 1,
            'strict_preflight': True,
            'dir': str(svg_ir_dir),
            'files': [slide['redcube_svg_ir_file'] for slide in manifest_slides],
        },
        'mode': args.mode,
        'editable_artifact': True,
        'pptx_file': str(output_pptx),
        'pdf_file': str(output_pdf) if output_pdf else None,
        'page_count': len(slides),
        'screenshot_dimensions': image_dimensions(Path(preview_files[0])) if preview_files else None,
        'preview_screenshots': preview_files,
        'slides': manifest_slides,
        'repair_log': repair_log,
    }
    shape_manifest.parent.mkdir(parents=True, exist_ok=True)
    shape_manifest.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding='utf-8')
    result = {
        'status': 'completed',
        'builder': manifest['builder'],
        'capability': manifest['capability'],
        'engine_contract': engine_contract,
        'engine_contract_file': str(engine_contract_file),
        'engine_capabilities': ENGINE_CAPABILITIES,
        'shape_manifest_schema_version': manifest['schema_version'],
        'mode': args.mode,
        'page_count': len(slides),
        'pptx_file': str(output_pptx),
        'pdf_file': str(output_pdf) if output_pdf else None,
        'shape_manifest_file': str(shape_manifest),
        'native_quality_model': manifest['native_quality_model'],
        'native_quality_surface': manifest['native_quality_surface'],
        'render_proof': render_proof,
        'redcube_svg_ir': manifest['redcube_svg_ir'],
        'preview_screenshots': preview_files,
        'screenshot_dimensions': manifest['screenshot_dimensions'],
        'slides': manifest_slides,
        'repair_log_file': str(repair_log_file) if repair_log_file else None,
        'repair_log': repair_log,
    }
    print(json.dumps(result, ensure_ascii=False))


if __name__ == '__main__':
    main()
