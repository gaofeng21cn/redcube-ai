#!/usr/bin/env python3
import argparse
import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def fail(message: str) -> None:
    print(message, file=sys.stderr)
    raise SystemExit(1)


def list_screenshots(screenshots_dir: Path):
    files = sorted(path for path in screenshots_dir.glob('*.png') if path.is_file())
    if not files:
        fail(f'No PNG screenshots found in {screenshots_dir}')
    return files


def build_pptx(images, output_pptx: Path):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    for image in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))


def build_pdf(images, output_pdf: Path):
    output_pdf.parent.mkdir(parents=True, exist_ok=True)
    pil_images = [Image.open(image).convert('RGB') for image in images]
    try:
        first, rest = pil_images[0], pil_images[1:]
        first.save(str(output_pdf), save_all=True, append_images=rest)
    finally:
        for image in pil_images:
            image.close()


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description='Export ppt_deck screenshots into PPTX/PDF artifacts.')
    parser.add_argument('--screenshots-dir', required=True)
    parser.add_argument('--output-pptx', required=True)
    parser.add_argument('--output-pdf')
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    screenshots_dir = Path(args.screenshots_dir).resolve()
    if not screenshots_dir.exists():
        fail(f'Screenshots directory not found: {screenshots_dir}')
    images = list_screenshots(screenshots_dir)

    output_pptx = Path(args.output_pptx).resolve()
    build_pptx(images, output_pptx)

    output_pdf = None
    if args.output_pdf:
        output_pdf = Path(args.output_pdf).resolve()
        build_pdf(images, output_pdf)

    result = {
        'status': 'completed',
        'converter': {
            'kind': 'python_playwright_pptx',
            'screenshots_dir': str(screenshots_dir),
        },
        'page_count': len(images),
        'pptx_file': str(output_pptx),
        'pdf_file': str(output_pdf) if output_pdf else None,
    }
    print(json.dumps(result, ensure_ascii=False))


if __name__ == '__main__':
    main()
