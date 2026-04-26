#!/usr/bin/env python3
import argparse
import json
import sys
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


CANVAS_PX = (1152, 648)
SLIDE_W = Inches(16)
SLIDE_H = Inches(9)
DEFAULT_ENGINE_CONTRACT_FILE = (
    Path(__file__).resolve().parents[3]
    / 'contracts'
    / 'runtime-program'
    / 'ppt-native-python-engine-contract.json'
)


def fail(message: str) -> None:
    print(message, file=sys.stderr)
    raise SystemExit(1)


def safe_text(value, fallback: str = '') -> str:
    text = str(value or '').strip()
    return text or fallback


def safe_list(value):
    return value if isinstance(value, list) else []


def load_engine_contract(contract_file: Path) -> dict:
    if not contract_file.exists():
        fail(f'engine contract not found: {contract_file}')
    contract = json.loads(contract_file.read_text(encoding='utf-8'))
    required = {
        'kind': 'redcube_native_ppt_python_engine',
        'language': 'python',
        'contract_version': 1,
        'input_boundary': 'slide_blueprint_plus_visual_direction_json',
        'review_boundary': 'rendered_pptx_screenshots',
    }
    for key, expected in required.items():
        if contract.get(key) != expected:
            fail(f'engine contract field mismatch: {key}')
    routes = contract.get('owned_routes')
    if routes != ['author_pptx_native', 'repair_pptx_native']:
        fail('engine contract owned_routes mismatch')
    return contract


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip('#')
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


PALETTES = [
    {'bg': '#F7F3EA', 'ink': '#18212B', 'muted': '#59636F', 'accent': '#BA3F1D', 'panel': '#EFE6D6'},
    {'bg': '#EEF4F1', 'ink': '#12231F', 'muted': '#52645F', 'accent': '#147C72', 'panel': '#DCEAE4'},
    {'bg': '#F3F1F7', 'ink': '#201B2C', 'muted': '#625A6F', 'accent': '#6E4AA8', 'panel': '#E4DFF0'},
    {'bg': '#F4F6F8', 'ink': '#152033', 'muted': '#5D6675', 'accent': '#2E6FBB', 'panel': '#E2E8F0'},
]


def font(size: int, bold: bool = False):
    candidates = [
        '/System/Library/Fonts/PingFang.ttc',
        '/System/Library/Fonts/Supplemental/Arial Unicode.ttf',
        '/Library/Fonts/Arial Unicode.ttf',
    ]
    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size=size)
        except Exception:
            continue
    return ImageFont.load_default()


def add_textbox(slide, left, top, width, height, text, size, color, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_panel(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(1)
    return shape


def slide_points(slide_data):
    points = [
        safe_text(item)
        for item in safe_list(slide_data.get('page_core_content'))
        if safe_text(item)
    ]
    if points:
        return points[:4]
    points = [
        safe_text(item)
        for item in safe_list(slide_data.get('evidence_points'))
        if safe_text(item)
    ]
    if points:
        return points[:4]
    return [safe_text(slide_data.get('core_sentence'), 'Key message')]


def draw_wrapped(draw, xy, text, image_font, fill, width_chars):
    lines = []
    for paragraph in safe_text(text).split('\n'):
        lines.extend(textwrap.wrap(paragraph, width=width_chars) or [''])
    x, y = xy
    line_height = int(image_font.size * 1.35) if hasattr(image_font, 'size') else 18
    for line in lines[:6]:
        draw.text((x, y), line, font=image_font, fill=fill)
        y += line_height


def render_preview(slide_data, index, total, preview_file: Path, palette, repaired: bool):
    image = Image.new('RGB', CANVAS_PX, palette['bg'])
    draw = ImageDraw.Draw(image)
    accent = palette['accent']
    ink = palette['ink']
    muted = palette['muted']
    panel = palette['panel']
    draw.rectangle([0, 0, 36, CANVAS_PX[1]], fill=accent)
    draw.rectangle([84, 116, 1068, 118], fill=accent)
    if repaired:
        draw.rectangle([1032, 36, 1084, 48], fill=accent)
    draw_wrapped(draw, (84, 54), safe_text(slide_data.get('title'), f'Slide {index}'), font(34, True), ink, 29)
    draw_wrapped(draw, (84, 142), safe_text(slide_data.get('core_sentence')), font(22), muted, 50)
    y = 250
    for point_index, point in enumerate(slide_points(slide_data)[:4], 1):
        draw.rounded_rectangle([96, y, 1056, y + 68], radius=18, fill=panel)
        draw.text((124, y + 18), f'{point_index}', font=font(24, True), fill=accent)
        draw_wrapped(draw, (176, y + 16), point, font(19), ink, 58)
        y += 86
    footer = f'{safe_text(slide_data.get("slide_id"), f"S{index:02d}")} / {index} of {total}'
    draw.text((84, 596), footer, font=font(16), fill=muted)
    preview_file.parent.mkdir(parents=True, exist_ok=True)
    image.save(preview_file)


def build_deck(slides, output_pptx: Path, preview_dir: Path, repaired_slide_ids):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]
    preview_files = []
    manifest_slides = []
    for index, slide_data in enumerate(slides, 1):
        palette = PALETTES[(index - 1) % len(PALETTES)]
        slide = prs.slides.add_slide(blank_layout)
        repaired = safe_text(slide_data.get('slide_id')) in repaired_slide_ids
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = rgb(palette['bg'])
        bg.line.color.rgb = rgb(palette['bg'])
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), SLIDE_H)
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(palette['accent'])
        bar.line.color.rgb = rgb(palette['accent'])
        add_textbox(
            slide,
            Inches(1.15),
            Inches(0.55),
            Inches(13.3),
            Inches(0.92),
            safe_text(slide_data.get('title'), f'Slide {index}'),
            32,
            rgb(palette['ink']),
            bold=True,
        )
        add_textbox(
            slide,
            Inches(1.15),
            Inches(1.62),
            Inches(13.3),
            Inches(0.72),
            safe_text(slide_data.get('core_sentence')),
            18,
            rgb(palette['muted']),
        )
        top = Inches(3.05)
        for point_index, point in enumerate(slide_points(slide_data)[:4], 1):
            panel = add_panel(slide, Inches(1.25), top, Inches(13.3), Inches(0.78), rgb(palette['panel']))
            add_textbox(
                slide,
                Inches(1.45),
                top + Inches(0.13),
                Inches(0.45),
                Inches(0.34),
                str(point_index),
                18,
                rgb(palette['accent']),
                bold=True,
            )
            add_textbox(
                slide,
                Inches(2.1),
                top + Inches(0.11),
                Inches(11.8),
                Inches(0.42),
                point,
                15,
                rgb(palette['ink']),
            )
            top += Inches(0.96)
        add_textbox(
            slide,
            Inches(1.15),
            Inches(8.25),
            Inches(13.3),
            Inches(0.3),
            f'{safe_text(slide_data.get("slide_id"), f"S{index:02d}")} / {index} of {len(slides)}',
            10,
            rgb(palette['muted']),
        )
        preview_file = preview_dir / f'{index:02d}-{safe_text(slide_data.get("slide_id"), f"S{index:02d}")}.png'
        render_preview(slide_data, index, len(slides), preview_file, palette, repaired)
        preview_files.append(str(preview_file))
        manifest_slides.append({
            'slide_id': safe_text(slide_data.get('slide_id'), f'S{index:02d}'),
            'title': safe_text(slide_data.get('title'), f'Slide {index}'),
            'layout_family': safe_text(slide_data.get('layout_family')),
            'shape_count': len(slide.shapes),
            'text_box_count': sum(1 for shape in slide.shapes if getattr(shape, 'has_text_frame', False)),
            'preview_screenshot_file': str(preview_file),
            'repaired': repaired,
        })
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))
    return preview_files, manifest_slides


def build_pdf(preview_files, output_pdf: Path):
    output_pdf.parent.mkdir(parents=True, exist_ok=True)
    images = [Image.open(file).convert('RGB') for file in preview_files]
    try:
        first, rest = images[0], images[1:]
        first.save(str(output_pdf), save_all=True, append_images=rest)
    finally:
        for image in images:
            image.close()


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description='Build editable native PPTX artifacts for ppt_deck.')
    parser.add_argument('--input-json', required=True)
    parser.add_argument('--mode', choices=['author', 'repair'], required=True)
    parser.add_argument('--output-pptx', required=True)
    parser.add_argument('--shape-manifest', required=True)
    parser.add_argument('--preview-dir', required=True)
    parser.add_argument('--output-pdf')
    parser.add_argument('--repair-log')
    parser.add_argument('--engine-contract', default=str(DEFAULT_ENGINE_CONTRACT_FILE))
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    payload = json.loads(Path(args.input_json).read_text(encoding='utf-8'))
    engine_contract_file = Path(args.engine_contract).resolve()
    engine_contract = load_engine_contract(engine_contract_file)
    blueprint = payload.get('blueprint') or {}
    slides = safe_list(blueprint.get('slides'))
    if not slides:
        fail('native PPT authoring requires slide_blueprint.slides')
    repair_feedback = safe_list(payload.get('repair_feedback'))
    repaired_slide_ids = {
        safe_text(item.get('slide_id'))
        for item in repair_feedback
        if isinstance(item, dict) and safe_text(item.get('slide_id'))
    }
    output_pptx = Path(args.output_pptx).resolve()
    shape_manifest = Path(args.shape_manifest).resolve()
    preview_dir = Path(args.preview_dir).resolve()
    preview_files, manifest_slides = build_deck(slides, output_pptx, preview_dir, repaired_slide_ids)
    output_pdf = Path(args.output_pdf).resolve() if args.output_pdf else None
    if output_pdf:
        build_pdf(preview_files, output_pdf)
    repair_log_file = Path(args.repair_log).resolve() if args.repair_log else None
    repair_log = {
        'mode': args.mode,
        'consumed_review_stage': 'screenshot_review' if args.mode == 'repair' else None,
        'target_slide_ids': sorted(repaired_slide_ids),
        'feedback_count': len(repair_feedback),
        'repair_log_file': str(repair_log_file) if repair_log_file else None,
    }
    if repair_log_file:
        repair_log_file.parent.mkdir(parents=True, exist_ok=True)
        repair_log_file.write_text(json.dumps(repair_log, ensure_ascii=False, indent=2), encoding='utf-8')
    manifest = {
        'schema_version': 1,
        'artifact_kind': 'ppt_deck_native_shape_manifest',
        'engine_contract': engine_contract,
        'engine_contract_file': str(engine_contract_file),
        'mode': args.mode,
        'editable_artifact': True,
        'pptx_file': str(output_pptx),
        'pdf_file': str(output_pdf) if output_pdf else None,
        'page_count': len(slides),
        'screenshot_dimensions': {'width': CANVAS_PX[0], 'height': CANVAS_PX[1]},
        'preview_screenshots': preview_files,
        'slides': manifest_slides,
        'repair_log': repair_log,
    }
    shape_manifest.parent.mkdir(parents=True, exist_ok=True)
    shape_manifest.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding='utf-8')
    result = {
        'status': 'completed',
        'builder': {'kind': 'python_pptx_native_shapes'},
        'engine_contract': engine_contract,
        'engine_contract_file': str(engine_contract_file),
        'shape_manifest_schema_version': manifest['schema_version'],
        'mode': args.mode,
        'page_count': len(slides),
        'pptx_file': str(output_pptx),
        'pdf_file': str(output_pdf) if output_pdf else None,
        'shape_manifest_file': str(shape_manifest),
        'preview_screenshots': preview_files,
        'screenshot_dimensions': manifest['screenshot_dimensions'],
        'slides': manifest_slides,
        'repair_log_file': str(repair_log_file) if repair_log_file else None,
        'repair_log': repair_log,
    }
    print(json.dumps(result, ensure_ascii=False))


if __name__ == '__main__':
    main()
